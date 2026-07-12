# -*- coding: utf-8 -*-
"""両家顔合わせ しおり → PowerPoint(.pptx) 生成スクリプト
HTML版と同じ内容・レイアウト（A4縦・全7ページ）を、PowerPointで
そのまま編集できる形で出力する。文字はすべてテキストボックス、
写真枠は差し替え用の図形として配置する。
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- パレット（HTML版と統一） ----
INK   = RGBColor(0x2E, 0x2A, 0x2B)
SOFT  = RGBColor(0x6B, 0x62, 0x60)
WINE  = RGBColor(0x7A, 0x2E, 0x3A)
GOLD  = RGBColor(0xB0, 0x8D, 0x4C)
GOLDS = RGBColor(0xD8, 0xC3, 0x9A)
FIELD = RGBColor(0xB9, 0x8A, 0x6B)
LINE  = RGBColor(0xE3, 0xD9, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MINCHO = "游明朝"
GOTHIC = "游ゴシック"

PW, PH = Cm(21.0), Cm(29.7)   # A4 縦

prs = Presentation()
prs.slide_width = PW
prs.slide_height = PH
BLANK = prs.slide_layouts[6]


def set_ea(run, name):
    """East Asian（日本語）書体を指定する。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def add_slide():
    s = prs.slides.add_slide(BLANK)
    # 背景は白（デフォルト）。金の外枠を描く。
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.0), Cm(1.0),
                               PW - Cm(2.0), PH - Cm(2.0))
    frame.fill.background()
    frame.line.color.rgb = GOLDS
    frame.line.width = Pt(0.75)
    frame.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, color, font, bold, spacing_pt, letter_kern), ...]
    1要素につき1段落。letter_kern は文字間(Pt)。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, color, font, bold, sp, kern = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sp is not None:
            p.space_after = Pt(sp)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        set_ea(r, font)
        if kern:
            r._r.get_or_add_rPr().set('spc', str(int(kern * 100)))
    return tb


def L(text, size, color=INK, font=GOTHIC, bold=False, sp=6, kern=0):
    return (text, size, color, font, bold, sp, kern)


def hline(slide, x, y, w, color=LINE, weight=0.75):
    ln = slide.shapes.add_connector(2, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def kicker(slide, no, label):
    txt(slide, Cm(2.0), Cm(1.9), Cm(10), Cm(0.8),
        [L(no + "   " + label, 11, WINE, GOTHIC, False, 0, 3)])
    hline(slide, Cm(2.0), Cm(2.75), Cm(17.0))


def heading(slide, jp, en, y=Cm(3.0)):
    txt(slide, Cm(2.0), y, Cm(17), Cm(1.3),
        [L(jp, 26, INK, MINCHO, True, 2, 2)])
    txt(slide, Cm(2.0), y + Cm(1.25), Cm(17), Cm(0.7),
        [L(en, 11, SOFT, GOTHIC, False, 0, 3)])


def photo_frame(slide, x, y, w, h):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.background()
    box.line.color.rgb = FIELD
    box.line.width = Pt(0.75)
    box.line.dash_style = 3  # dash
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "写真"
    r.font.size = Pt(10)
    r.font.color.rgb = FIELD
    r.font.name = GOTHIC
    set_ea(r, GOTHIC)
    return box


# =============== 1. 表紙 ===============
s = add_slide()
cx, cw = Cm(2.0), Cm(17.0)
txt(s, cx, Cm(5.2), cw, Cm(0.8), [L("結び", 12, WINE, MINCHO, False, 0, 8)], PP_ALIGN.CENTER)
txt(s, cx, Cm(8.6), cw, Cm(2.0), [L("両家顔合わせ", 40, INK, MINCHO, True, 0, 8)], PP_ALIGN.CENTER)
txt(s, cx, Cm(11.0), cw, Cm(0.8), [L("FAMILY  INTRODUCTION", 12, GOLD, GOTHIC, False, 0, 6)], PP_ALIGN.CENTER)
txt(s, cx, Cm(14.2), cw, Cm(1.2),
    [L("新郎 名前　×　新婦 名前", 20, INK, MINCHO, False, 0, 3)], PP_ALIGN.CENTER)
txt(s, cx, Cm(15.7), cw, Cm(0.8),
    [L("〇〇家　・　〇〇家", 13, SOFT, GOTHIC, False, 0, 4)], PP_ALIGN.CENTER)
txt(s, cx, Cm(19.2), cw, Cm(3.0), [
    L("日　時　　20〇〇年〇月〇日（〇）　午前・午後〇時", 13, INK, GOTHIC, False, 10, 2),
    L("会　場　　会場名・お部屋", 13, INK, GOTHIC, False, 0, 2),
], PP_ALIGN.CENTER)

# =============== 2. ご挨拶 ＋ 当日の流れ ===============
s = add_slide()
kicker(s, "01", "ごあいさつ")
heading(s, "ご挨拶", "GREETING")
txt(s, Cm(2.0), Cm(5.6), Cm(17.0), Cm(6.0), [
    L("本日はお忙しい中、両家の顔合わせにお集まりいただき、誠にありがとうございます。", 13, INK, GOTHIC, False, 10),
    L("（ここに、お集まりへの感謝と、この日を迎えた喜びを書きます。）", 13, FIELD, GOTHIC, True, 16),
    L("本日はどうぞお時間の許す限り、和やかにお過ごしいただけましたら幸いです。これからの両家のお付き合いも、末永くよろしくお願い申し上げます。", 13, INK, GOTHIC, False, 0),
], anchor=MSO_ANCHOR.TOP)
# 当日の流れ
txt(s, Cm(2.0), Cm(13.2), Cm(10), Cm(0.9),
    [L("当日の流れ　　FLOW OF THE DAY", 15, WINE, MINCHO, False, 0, 2)])
hline(s, Cm(2.0), Cm(14.15), Cm(17.0))
txt(s, Cm(2.0), Cm(14.7), Cm(17.0), Cm(2.2),
    [L("始まりの挨拶　›　両家の紹介　›　観覧　›　食事　›　記念撮影　›　結びの挨拶",
       16, INK, GOTHIC, False, 0, 1)])

# =============== 3. 二人のプロフィール ===============
s = add_slide()
kicker(s, "02", "プロフィール")
heading(s, "二人のプロフィール", "PROFILE")
rows = ["ふりがな", "生年月日", "出身地", "血液型", "お仕事", "趣味", "ひとこと"]
vals = ["ふりがな", "〇年〇月〇日", "〇〇県", "〇型", "職業", "趣味・好きなこと", "自己紹介・意気込み"]
def profile_col(x, role, name):
    txt(s, x, Cm(5.4), Cm(7.5), Cm(0.7), [L(role, 12, GOLD, GOTHIC, False, 0, 4)])
    txt(s, x, Cm(6.0), Cm(7.5), Cm(1.0), [L(name, 18, INK, MINCHO, False, 0, 2)])
    y = Cm(7.5)
    lines = []
    for k, v in zip(rows, vals):
        lines.append((k + "：" + v, 13, INK, GOTHIC, False, 12, 0))
    tb = s.shapes.add_textbox(x, y, Cm(7.6), Cm(14))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_top=0; tf.margin_right=0; tf.margin_bottom=0
    for i,(k,v) in enumerate(zip(rows, vals)):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(11)
        rk = p.add_run(); rk.text = k + "　"
        rk.font.size=Pt(12); rk.font.color.rgb=SOFT; rk.font.name=GOTHIC; set_ea(rk,GOTHIC)
        rv = p.add_run(); rv.text = v
        rv.font.size=Pt(13); rv.font.color.rgb=FIELD; rv.font.bold=True; rv.font.name=GOTHIC; set_ea(rv,GOTHIC)
profile_col(Cm(2.0), "GROOM", "新郎 氏名")
profile_col(Cm(11.4), "BRIDE", "新婦 氏名")
# 中央の縦線
vln = s.shapes.add_connector(2, Cm(10.5), Cm(5.6), Cm(10.5), Cm(20.5))
vln.line.color.rgb = LINE; vln.line.width = Pt(0.75)

# =============== 4. 両家の家族紹介 ===============
s = add_slide()
kicker(s, "03", "家族紹介")
heading(s, "両家の家族紹介", "OUR FAMILIES")
def family(x0, house, members):
    txt(s, x0, Cm(5.5), Cm(7.6), Cm(0.8), [L(house, 14, WINE, MINCHO, False, 0, 2)])
    hline(s, x0, Cm(6.35), Cm(7.6))
    pw, ph = Cm(3.5), Cm(4.4)
    gx, gy = Cm(0.6), Cm(0.5)
    for i, rel in enumerate(members):
        col = i % 2; row = i // 2
        x = x0 + col * (pw + gx)
        y = Cm(6.8) + row * (ph + Cm(2.4))
        photo_frame(s, x, y, pw, ph)
        txt(s, x, y + ph + Cm(0.1), pw, Cm(1.8), [
            L(rel, 11, GOLD, GOTHIC, False, 2, 3),
            L("お名前", 13, FIELD, GOTHIC, True, 2),
            L("ひとこと", 11, SOFT, GOTHIC, False, 0),
        ], PP_ALIGN.CENTER)
family(Cm(2.0), "〇〇家（新郎側）", ["父", "母", "姉", "兄"])
family(Cm(11.4), "〇〇家（新婦側）", ["父", "母", "弟"])

# =============== 5. 二人の思い出 ===============
s = add_slide()
kicker(s, "04", "ストーリー")
heading(s, "二人の思い出", "FROM OUR FIRST MEETING TO THE PROPOSAL")
memories = [
    ("20〇〇年〇月　—　出会い", "はじめて出会った場所・きっかけ", "そのときの第一印象や、どんな出会いだったかを書きます。"),
    ("20〇〇年〇月　—　交際スタート", "お付き合いが始まった日", "告白のエピソードや、初デートの思い出など。"),
    ("20〇〇年〇月　—　思い出の出来事", "旅行・記念日など心に残る出来事", "二人の絆が深まったエピソードを。"),
    ("20〇〇年〇月　—　プロポーズ", "プロポーズの言葉・場所", "どんなプロポーズだったか、その時の気持ちを書きます。"),
]
y = Cm(5.6)
for when, what, desc in memories:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.0), y + Cm(0.15), Cm(0.28), Cm(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = WINE; dot.line.fill.background(); dot.shadow.inherit=False
    txt(s, Cm(2.7), y, Cm(16.3), Cm(3.6), [
        L(when, 12, GOLD, GOTHIC, False, 3, 2),
        L(what, 15, INK, MINCHO, False, 3, 1),
        L(desc, 12, SOFT, GOTHIC, False, 0),
    ])
    y += Cm(4.05)

# =============== 6. 今後の予定 ===============
s = add_slide()
kicker(s, "05", "これから")
heading(s, "今後の予定", "OUR PLANS")
plans = [
    ("入籍", "入籍", "20〇〇年〇月〇日 に入籍予定です。（記念日にちなんだ日など）"),
    ("新居", "新生活・新居", "〇〇（エリア）で新生活をスタートします。（時期・住まいのことなど）"),
    ("挙式", "結婚式", "20〇〇年〇月頃、〇〇（会場）にて挙式予定です。（未定の場合は「検討中」でも）"),
    ("他", "これからのこと", "新婚旅行や、これからの二人の目標など、自由に書けます。"),
]
y = Cm(5.8)
for badge, title, body in plans:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2.0), y, Cm(17.0), Cm(3.4))
    card.fill.background(); card.line.color.rgb = LINE; card.line.width = Pt(0.75); card.shadow.inherit=False
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.7), y + Cm(0.9), Cm(1.6), Cm(1.6))
    circ.fill.background(); circ.line.color.rgb = GOLDS; circ.line.width = Pt(1.0); circ.shadow.inherit=False
    ctf = circ.text_frame; cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = badge; cr.font.size=Pt(13); cr.font.color.rgb=WINE; cr.font.name=MINCHO; set_ea(cr,MINCHO)
    txt(s, Cm(5.0), y + Cm(0.6), Cm(13.6), Cm(2.4), [
        L(title, 16, INK, MINCHO, False, 4, 1),
        L(body, 12, SOFT, GOTHIC, False, 0),
    ], anchor=MSO_ANCHOR.MIDDLE)
    y += Cm(3.9)

# =============== 7. 感謝のメッセージ ===============
s = add_slide()
txt(s, Cm(2.0), Cm(6.0), Cm(17.0), Cm(1.2), [L("結", 30, GOLD, MINCHO, False, 0)], PP_ALIGN.CENTER)
txt(s, Cm(2.0), Cm(8.4), Cm(17.0), Cm(1.3), [L("感謝をこめて", 26, INK, MINCHO, True, 2, 2)], PP_ALIGN.CENTER)
txt(s, Cm(2.0), Cm(10.0), Cm(17.0), Cm(0.7), [L("WITH GRATITUDE", 11, SOFT, GOTHIC, False, 0, 3)], PP_ALIGN.CENTER)
txt(s, Cm(3.0), Cm(12.5), Cm(15.0), Cm(6.0), [
    L("（ここに、両家のご家族へ向けた感謝の言葉を書きます。", 15, FIELD, MINCHO, True, 8),
    L("これまで育ててくださったことへの感謝と、", 15, FIELD, MINCHO, True, 8),
    L("これから家族として歩んでいく決意などを綴ります。）", 15, FIELD, MINCHO, True, 0),
], PP_ALIGN.CENTER)
txt(s, Cm(2.0), Cm(20.0), Cm(17.0), Cm(1.0),
    [L("新郎 名前　・　新婦 名前", 15, INK, MINCHO, False, 0, 3)], PP_ALIGN.CENTER)
txt(s, Cm(2.0), Cm(21.0), Cm(17.0), Cm(0.8),
    [L("20〇〇年〇月〇日", 12, SOFT, GOTHIC, False, 0, 3)], PP_ALIGN.CENTER)
txt(s, Cm(2.0), Cm(26.6), Cm(17.0), Cm(0.8),
    [L("本日は誠にありがとうございました", 11, SOFT, MINCHO, False, 0, 4)], PP_ALIGN.CENTER)

out = "/home/user/simple-alarm-app/wedding-booklet/kaoawase-shiori.pptx"
prs.save(out)
print("saved", out)
